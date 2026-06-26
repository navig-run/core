"""
navig.inbox.extract — universal content extraction for the inbox.

Turns ANY dropped file (text, image, audio, video, pdf, office doc, or unknown
binary) into normalized text + metadata, so the markdown-only inbox classifier
can route it.  Nothing is ever dropped: an unreadable binary still yields a stub
document so the file is preserved and visible.

Design rules
------------
* **Local-first.** The common path (text / OCR / pdf / transcription) reuses
  standalone helpers (``navig.core.ocr``, ``navig.voice.stt``, ``pypdf``) and the
  stdlib only — it never imports the heavy ``navig.gateway.*`` cascade.
* **Cloud is opt-in.** GPT-4o vision (and cloud STT) run only when the resolved
  :class:`ExtractPolicy` allows it; those stages (and the ``BudgetGuard``) are
  lazy-imported solely in that branch.
* **Never raises.** Every optional dependency (pytesseract, ffmpeg, pypdf,
  PyMuPDF, whisper, python-docx, openpyxl, mutagen, Pillow) degrades to
  "metadata only" with an ``errors[]`` entry instead of crashing the pipeline.

Public API
----------
``extract(path, *, policy=None, budget=None, cache=None) -> ExtractResult``
``to_markdown(result, *, source_label, original_preserved) -> str``
``ExtractPolicy`` / ``ExtractResult``
"""

from __future__ import annotations

import asyncio
import logging
import mimetypes
import subprocess
import tempfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger("navig.inbox.extract")

# Bumped when extraction output format changes, to invalidate stale cache entries.
EXTRACT_VERSION = 1

# Max raw bytes we will read into memory for a single non-text file.
_DEFAULT_MAX_BYTES = 50_000_000  # 50 MB

# ── Suffix → kind dispatch ────────────────────────────────────────────────────

_TEXT_SUFFIXES = {
    ".md", ".markdown", ".txt", ".rst", ".text", ".log",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".csv", ".tsv",
    ".py", ".ts", ".tsx", ".js", ".jsx", ".mjs", ".cjs", ".rs", ".go",
    ".java", ".kt", ".c", ".h", ".cpp", ".hpp", ".cs", ".rb", ".php",
    ".sh", ".ps1", ".sql", ".html", ".htm", ".css", ".scss", ".xml", ".svg",
}
_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff", ".heic"}
_AUDIO_SUFFIXES = {".mp3", ".wav", ".m4a", ".ogg", ".oga", ".flac", ".aac", ".wma", ".opus"}
_VIDEO_SUFFIXES = {".mp4", ".mov", ".mkv", ".webm", ".avi", ".m4v"}
_PDF_SUFFIXES = {".pdf"}
_DOCX_SUFFIXES = {".docx"}
_XLSX_SUFFIXES = {".xlsx"}
_PPTX_SUFFIXES = {".pptx"}


# ── Policy & result ───────────────────────────────────────────────────────────


@dataclass
class ExtractPolicy:
    """How aggressively to extract. ``mode`` is selectable in the deck."""

    mode: str = "auto"          # "local" | "auto" | "cloud"
    max_pdf_pages: int = 30
    max_bytes: int = _DEFAULT_MAX_BYTES

    @property
    def allow_cloud(self) -> bool:
        """auto/cloud may call paid providers (still key- and budget-gated)."""
        return self.mode != "local"

    @property
    def prefer_local_stt(self) -> bool:
        """local/auto transcribe with local Whisper first."""
        return self.mode != "cloud"

    @classmethod
    def from_config(cls) -> ExtractPolicy:
        mode = "auto"
        try:
            from navig.config import get_config_manager

            cfg = get_config_manager().global_config or {}
            raw = ((cfg.get("inbox") or {}).get("extract") or {}).get("mode")
            if isinstance(raw, str) and raw.strip().lower() in {"local", "auto", "cloud"}:
                mode = raw.strip().lower()
        except Exception as exc:  # noqa: BLE001
            logger.debug("ExtractPolicy.from_config fell back to auto: %s", exc)
        return cls(mode=mode)


@dataclass
class ExtractResult:
    text: str = ""
    kind: str = "binary"        # text|image|audio|video|pdf|docx|xlsx|pptx|binary
    metadata: dict[str, Any] = field(default_factory=dict)
    extracted_by: list[str] = field(default_factory=list)
    original_path: str = ""
    content_hash: str = ""
    cached: bool = False
    errors: list[str] = field(default_factory=list)

    # Subset persisted to the content-addressed cache (machine-independent).
    _CACHE_KEYS = ("text", "kind", "metadata", "extracted_by", "errors")

    def cacheable(self) -> dict[str, Any]:
        d = asdict(self)
        return {"_v": EXTRACT_VERSION, **{k: d[k] for k in self._CACHE_KEYS}}


# ── Public entry point ────────────────────────────────────────────────────────


def extract(
    path: str | Path,
    *,
    policy: ExtractPolicy | None = None,
    budget: Any | None = None,
    cache: Any | None = None,
) -> ExtractResult:
    """Extract normalized text + metadata from *path*. Never raises."""
    p = Path(path)
    policy = policy or ExtractPolicy.from_config()
    res = ExtractResult(original_path=str(p))

    try:
        size = p.stat().st_size
    except OSError as exc:
        res.errors.append(f"stat failed: {exc}")
        return res

    kind = _kind_for(p)
    res.kind = kind

    # Text fast-path: read directly, no byte cap, no cache, no gateway.
    if kind == "text":
        try:
            res.text = p.read_text(encoding="utf-8", errors="replace")
            res.extracted_by = ["passthrough"]
        except OSError as exc:
            res.errors.append(f"read failed: {exc}")
        return res

    # Binary-ish: read bytes once (capped), hash, consult cache.
    if size > policy.max_bytes:
        res.metadata = {"size_bytes": size, "skipped": "too_large"}
        res.errors.append(f"file too large ({size} bytes > {policy.max_bytes})")
        return res
    try:
        raw = p.read_bytes()
    except OSError as exc:
        res.errors.append(f"read failed: {exc}")
        return res

    res.content_hash = _sha256(raw)
    if cache is not None:
        hit = cache.get(res.content_hash)
        if isinstance(hit, dict) and hit.get("_v") == EXTRACT_VERSION:
            for k in ExtractResult._CACHE_KEYS:
                if k in hit:
                    setattr(res, k, hit[k])
            res.cached = True
            return res

    dispatch = {
        "image": _extract_image,
        "audio": _extract_audio,
        "video": _extract_video,
        "pdf": _extract_pdf_doc,
        "docx": _extract_docx,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
    }.get(kind)

    if dispatch is None:
        res.metadata = {"size_bytes": size, "mime": mimetypes.guess_type(str(p))[0] or ""}
        res.errors.append("unsupported file type — preserved as-is")
    else:
        try:
            dispatch(p, raw, policy, budget, res)
        except Exception as exc:  # noqa: BLE001 — extraction must never crash routing
            logger.warning("extract(%s) stage failed: %s", p.name, exc)
            res.errors.append(f"{kind} extraction failed: {exc}")

    if cache is not None and res.content_hash:
        try:
            cache.put(res.content_hash, res.cacheable())
        except Exception:  # noqa: BLE001
            pass  # best-effort
    return res


def to_markdown(res: ExtractResult, *, source_label: str, original_preserved: str = "") -> str:
    """Render an ExtractResult as a normalized markdown doc the Classifier reads."""
    fm: list[str] = ["---", f"source: {source_label}", f"kind: {res.kind}"]
    if res.extracted_by:
        fm.append(f"extracted_by: [{', '.join(res.extracted_by)}]")
    if res.content_hash:
        fm.append(f"content_hash: {res.content_hash}")
    if original_preserved:
        fm.append(f"original_preserved: {original_preserved}")
    gps = (res.metadata or {}).get("gps")
    if gps:
        fm.append(f"gps: {gps}")
    for key in ("title", "artist", "album", "camera", "taken_at", "duration_sec", "pages"):
        val = (res.metadata or {}).get(key)
        if val:
            fm.append(f"{key}: {val}")
    if res.errors:
        fm.append(f"extract_errors: [{'; '.join(res.errors)}]")
    fm.append("---")

    title = Path(source_label).stem.replace("-", " ").replace("_", " ").strip() or source_label
    body = res.text.strip() or f"_(no text extracted from {res.kind} file — original preserved)_"
    return f"{chr(10).join(fm)}\n\n# {title}\n\n{body}\n"


# ── Kind detection ────────────────────────────────────────────────────────────


def _kind_for(p: Path) -> str:
    suf = p.suffix.lower()
    if suf in _TEXT_SUFFIXES:
        return "text"
    if suf in _IMAGE_SUFFIXES:
        return "image"
    if suf in _AUDIO_SUFFIXES:
        return "audio"
    if suf in _VIDEO_SUFFIXES:
        return "video"
    if suf in _PDF_SUFFIXES:
        return "pdf"
    if suf in _DOCX_SUFFIXES:
        return "docx"
    if suf in _XLSX_SUFFIXES:
        return "xlsx"
    if suf in _PPTX_SUFFIXES:
        return "pptx"
    mime = (mimetypes.guess_type(str(p))[0] or "").lower()
    if mime.startswith("text/") or mime in {"application/json", "application/xml"}:
        return "text"
    if mime.startswith("image/"):
        return "image"
    if mime.startswith("audio/"):
        return "audio"
    if mime.startswith("video/"):
        return "video"
    if mime == "application/pdf":
        return "pdf"
    return "binary"


# ── Stage implementations (each mutates ``res``; never raises) ─────────────────


def _extract_image(p: Path, raw: bytes, policy: ExtractPolicy, budget: Any, res: ExtractResult) -> None:
    res.metadata.update(_image_metadata(raw))
    res.extracted_by.append("pillow")

    from navig.core.ocr import extract_ocr_text_from_image_bytes

    ocr = extract_ocr_text_from_image_bytes(raw)
    parts: list[str] = []
    if ocr:
        parts.append(ocr)
        res.extracted_by.append("tesseract")

    if policy.allow_cloud:
        desc = _maybe_vision(raw, budget, res)
        if desc:
            parts.append(f"**Visual description:** {desc}")
            res.extracted_by.append("vision")
    res.text = "\n\n".join(parts)


def _extract_audio(p: Path, raw: bytes, policy: ExtractPolicy, budget: Any, res: ExtractResult) -> None:
    res.metadata.update(_audio_metadata(raw))
    res.extracted_by.append("mutagen")
    transcript, provider = _transcribe(p, policy, res)
    if transcript:
        res.text = transcript
        res.extracted_by.append(provider)


def _extract_video(p: Path, raw: bytes, policy: ExtractPolicy, budget: Any, res: ExtractResult) -> None:
    import shutil

    ffmpeg = shutil.which("ffmpeg")
    if not ffmpeg:
        res.errors.append("ffmpeg not found — video audio-track not transcribed")
        return
    tmp = Path(tempfile.gettempdir()) / f"navig-extract-{res.content_hash[:8]}.wav"
    try:
        subprocess.run(
            [ffmpeg, "-y", "-i", str(p), "-vn", "-ac", "1", "-ar", "16000", str(tmp)],
            capture_output=True, timeout=600, check=False,
        )
        if tmp.exists() and tmp.stat().st_size > 0:
            transcript, provider = _transcribe(tmp, policy, res)
            if transcript:
                res.text = transcript
                res.extracted_by.append(provider)
        else:
            res.errors.append("ffmpeg produced no audio track")
    finally:
        try:
            tmp.unlink(missing_ok=True)
        except OSError:
            pass


def _extract_pdf_doc(p: Path, raw: bytes, policy: ExtractPolicy, budget: Any, res: ExtractResult) -> None:
    text, meta = _extract_pdf(p, max_pages=policy.max_pdf_pages, errors=res.errors)
    res.text = text
    res.metadata.update(meta)
    res.extracted_by.append("ocr" if meta.get("ocr_fallback") else "pypdf")


def _extract_docx(p: Path, raw: bytes, policy: ExtractPolicy, budget: Any, res: ExtractResult) -> None:
    try:
        import docx  # type: ignore

        doc = docx.Document(str(p))
        res.text = "\n".join(par.text for par in doc.paragraphs if par.text).strip()
        res.extracted_by.append("python-docx")
    except ImportError:
        res.errors.append("python-docx not installed — .docx not extracted")
    except Exception as exc:  # noqa: BLE001
        res.errors.append(f"docx parse failed: {exc}")


def _extract_xlsx(p: Path, raw: bytes, policy: ExtractPolicy, budget: Any, res: ExtractResult) -> None:
    try:
        import openpyxl  # type: ignore

        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        lines: list[str] = []
        for ws in wb.worksheets:
            lines.append(f"## {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append(" | ".join(cells))
        res.text = "\n".join(lines).strip()
        res.extracted_by.append("openpyxl")
    except ImportError:
        res.errors.append("openpyxl not installed — .xlsx not extracted")
    except Exception as exc:  # noqa: BLE001
        res.errors.append(f"xlsx parse failed: {exc}")


def _extract_pptx(p: Path, raw: bytes, policy: ExtractPolicy, budget: Any, res: ExtractResult) -> None:
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(p))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
                    lines.append(shape.text_frame.text)
        res.text = "\n".join(lines).strip()
        res.extracted_by.append("python-pptx")
    except ImportError:
        res.errors.append("python-pptx not installed — .pptx not extracted")
    except Exception as exc:  # noqa: BLE001
        res.errors.append(f"pptx parse failed: {exc}")


# ── Helpers (standalone — no gateway import) ──────────────────────────────────


def _sha256(data: bytes) -> str:
    import hashlib

    return hashlib.sha256(data).hexdigest()


def _image_metadata(raw: bytes) -> dict[str, Any]:
    """Minimal Pillow metadata (dims/format + best-effort camera/GPS). Standalone."""
    out: dict[str, Any] = {}
    try:
        import io

        from PIL import Image  # type: ignore

        img = Image.open(io.BytesIO(raw))
        out["width"], out["height"] = img.size
        out["format"] = img.format or "?"
        try:
            from PIL.ExifTags import GPSTAGS, TAGS  # type: ignore

            raw_exif = img._getexif() if hasattr(img, "_getexif") else None
            if raw_exif:
                named, gps_raw = {}, {}
                for tag_id, value in raw_exif.items():
                    tag = TAGS.get(tag_id, tag_id)
                    if tag == "GPSInfo" and isinstance(value, dict):
                        for g_id, g_val in value.items():
                            gps_raw[GPSTAGS.get(g_id, g_id)] = g_val
                    else:
                        try:
                            named[tag] = str(value)[:200]
                        except Exception:  # noqa: BLE001
                            pass
                out["camera"] = named.get("Model") or named.get("Make")
                out["taken_at"] = named.get("DateTimeOriginal") or named.get("DateTime")
                if gps_raw.get("GPSLatitude") and gps_raw.get("GPSLongitude"):
                    def _deg(vals, ref):
                        d, m, s = (float(v) for v in vals)
                        deg = d + m / 60 + s / 3600
                        return round(-deg if ref in ("S", "W") else deg, 6)

                    try:
                        out["gps"] = {
                            "lat": _deg(gps_raw["GPSLatitude"], gps_raw.get("GPSLatitudeRef", "N")),
                            "lon": _deg(gps_raw["GPSLongitude"], gps_raw.get("GPSLongitudeRef", "E")),
                        }
                    except Exception:  # noqa: BLE001
                        pass
        except Exception:  # noqa: BLE001
            pass  # EXIF best-effort
    except Exception as exc:  # noqa: BLE001
        logger.debug("image metadata: %s", exc)
    return {k: v for k, v in out.items() if v}


def _audio_metadata(raw: bytes) -> dict[str, Any]:
    """Minimal mutagen metadata (duration/title/artist/album). Standalone."""
    out: dict[str, Any] = {}
    try:
        import io

        import mutagen  # type: ignore

        f = mutagen.File(io.BytesIO(raw))
        if f is None:
            return out
        dur = getattr(getattr(f, "info", None), "length", None)
        if dur:
            out["duration_sec"] = int(dur)
        tags = f.tags or {}

        def _tag(*keys):
            for k in keys:
                v = tags.get(k)
                if v:
                    return str(v[0]) if isinstance(v, (list, tuple)) else str(v)
            return None

        out["title"] = _tag("TIT2", "\xa9nam", "title")
        out["artist"] = _tag("TPE1", "\xa9ART", "artist")
        out["album"] = _tag("TALB", "\xa9alb", "album")
    except Exception as exc:  # noqa: BLE001
        logger.debug("audio metadata: %s", exc)
    return {k: v for k, v in out.items() if v}


def _extract_pdf(p: Path, *, max_pages: int, errors: list[str]) -> tuple[str, dict[str, Any]]:
    """pypdf text; OCR-rasterize scanned PDFs via PyMuPDF + Tesseract. Standalone."""
    text_parts: list[str] = []
    pages = 0
    try:
        import pypdf  # type: ignore

        reader = pypdf.PdfReader(str(p))
        for i, page in enumerate(reader.pages):
            if i >= max_pages:
                break
            try:
                text_parts.append(page.extract_text() or "")
            except Exception:  # noqa: BLE001
                text_parts.append("")
            pages += 1
    except ImportError:
        errors.append("pypdf not installed — PDF text not extracted")
    except Exception as exc:  # noqa: BLE001
        errors.append(f"pypdf parse failed: {exc}")

    text = "\n".join(text_parts).strip()
    ocr_fallback = False
    if len(text) < 20:  # scanned / image-only PDF → rasterize + OCR
        try:
            import fitz  # type: ignore  # PyMuPDF

            from navig.core.ocr import extract_ocr_text_from_image_bytes

            doc = fitz.open(str(p))
            ocr_parts: list[str] = []
            for i, page in enumerate(doc):
                if i >= max_pages:
                    break
                pix = page.get_pixmap(dpi=150)
                t = extract_ocr_text_from_image_bytes(pix.tobytes("png"))
                if t:
                    ocr_parts.append(t)
            doc.close()
            if ocr_parts:
                text = "\n".join(ocr_parts).strip()
                ocr_fallback = True
                pages = pages or len(ocr_parts)
        except ImportError:
            errors.append("PyMuPDF not installed — scanned PDF not OCR'd")
        except Exception as exc:  # noqa: BLE001
            errors.append(f"PDF OCR fallback failed: {exc}")
    return text, {"pages": pages, "ocr_fallback": ocr_fallback}


def _transcribe(path: Path, policy: ExtractPolicy, res: ExtractResult) -> tuple[str | None, str]:
    """Transcribe audio via navig.voice.stt (standalone). Local-first per policy."""
    try:
        from navig.voice.stt import STT, STTConfig, STTProvider

        if policy.mode == "local":
            cfg = STTConfig(provider=STTProvider.WHISPER_LOCAL, fallback_providers=[])
        elif policy.mode == "cloud":
            cfg = STTConfig()  # WHISPER_API first, local fallback (lib default)
        else:  # auto: local-first, cloud fallback
            cfg = STTConfig(
                provider=STTProvider.WHISPER_LOCAL,
                fallback_providers=[STTProvider.WHISPER_API],
            )
        result = _run_async(STT(cfg).transcribe(str(path)))
        if result and getattr(result, "success", False) and result.text:
            prov = result.provider.value if getattr(result, "provider", None) else "stt"
            return result.text.strip(), prov
        if result and getattr(result, "error", None):
            res.errors.append(f"transcription: {result.error}")
    except Exception as exc:  # noqa: BLE001
        res.errors.append(f"transcription failed: {exc}")
    return None, "stt"


def _maybe_vision(raw: bytes, budget: Any, res: ExtractResult) -> str | None:
    """Cloud GPT-4o describe — lazy-imports the gateway-coupled stage + budget."""
    try:
        from navig.gateway.channels.media_engine.budget import BudgetGuard
        from navig.gateway.channels.media_engine.image import _stage_vision

        guard = budget or BudgetGuard()
        return _run_async(_stage_vision(raw, guard))
    except Exception as exc:  # noqa: BLE001
        res.errors.append(f"vision describe skipped: {exc}")
        return None


def _run_async(coro: Any) -> Any:
    """Run *coro* to completion whether or not an event loop is already running."""
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        return asyncio.run(coro)
    # A loop is running on this thread — execute in a fresh worker thread.
    import concurrent.futures

    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as ex:
        return ex.submit(lambda: asyncio.run(coro)).result()
